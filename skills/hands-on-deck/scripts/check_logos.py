#!/usr/bin/env python3
"""check_logos.py — post-build visual gate that catches the doubled-logo defect.

THE BUG THIS EXISTS FOR
-----------------------
The bespoke-deck pipeline places full-bleed images: Route 1 renders one PNG per
slide and drops it edge-to-edge; Route 2 adds full-bleed pictures via `add-picture`.
When a picture that has the EVERY logo / footer baked into its pixels lands on a
master or layout that ALREADY carries the logo/footer, you get two logos stacked.

This is invisible to image-only QA — the reviewer looks at the rendered PNG and
never renders the master *underneath* it, so the pass feels clean while the shipped
slide shows a doubled logo. Natalia and Mike both flagged it independently (2026-07-20),
and it has recurred several times.

THE RULE (the fix)
------------------
A full-bleed / footer-covering picture must leave the logo+footer zone EMPTY and let
the master supply the branding. Two sanctioned ways:
  (a) build the image with no logo and a reserved bottom margin, then drop it on the
      branded layout; or
  (b) drop it on a BLANK layout (slide_layouts[6]) and let the picture own the whole
      frame — never a branded layout + a logo-bearing picture.

WHAT THIS GATE DOES
-------------------
Structural, not pixel-based (so it's deterministic, fast, and CI-safe — no rendering):
for every slide it asks "does a slide-level picture cover the footer/logo zone that
the slide's layout or master ALSO fills with real branding?" and flags the slide if so.

  - "branding in the footer zone" = a non-placeholder graphic/logo on the layout or
    master whose box sits in the bottom band (backgrounds that cover the whole slide
    are ignored — they're not logos), or a layout placeholder explicitly named
    logo/footer.
  - A picture is a problem only if it covers >= --coverage of that zone. A picture
    that reserves the bottom margin (sanctioned fix a) covers ~0 of the zone and
    passes. A full-bleed picture on a blank, unbranded layout (sanctioned fix b) has
    no branding to collide with and passes.

Exit codes: 0 = clean, 1 = at least one slide flagged, 2 = usage/error.
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _rect(shape):
    """Absolute EMU box (l, t, w, h) or None if geometry is inherited/unset."""
    l, t, w, h = shape.left, shape.top, shape.width, shape.height
    if None in (l, t, w, h):
        return None
    return (int(l), int(t), int(w), int(h))


def _intersect_area(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix, iy = max(ax, bx), max(ay, by)
    ix2, iy2 = min(ax + aw, bx + bw), min(ay + ah, by + bh)
    if ix2 <= ix or iy2 <= iy:
        return 0
    return (ix2 - ix) * (iy2 - iy)


def _is_placeholder(shape):
    try:
        return bool(shape.is_placeholder)
    except Exception:
        return False


def footer_branding(layout, slide_w, slide_h, band):
    """Return (zone_rect, [(source, name, rect), ...]) for real branding in the
    bottom band that would render *under* a slide-level full-bleed picture."""
    zone = (0, int(slide_h * (1 - band)), slide_w, int(slide_h * band))
    slide_area = slide_w * slide_h
    found = []

    def consider(shape, source):
        r = _rect(shape)
        if not r:
            return
        if r[2] * r[3] > 0.7 * slide_area:  # a background fill, not a logo
            return
        if _is_placeholder(shape):
            # Non-placeholder graphics (a real logo picture/shape the designer added
            # to the footer zone) always count — handled below. Placeholders are the
            # tricky part: the template ships date / slide-number / body / EMPTY footer
            # placeholders on every master and layout, all with default field or prompt
            # text, and none of them are branding. A placeholder counts only when it is
            # a *logo* placeholder, or a *footer* placeholder that actually carries text.
            name = (shape.name or "").lower()
            if "logo" in name:
                pass  # a logo placeholder is branding regardless of text
            elif "footer" in name:
                has_text = False
                try:
                    has_text = bool(shape.has_text_frame and shape.text_frame.text.strip())
                except Exception:
                    has_text = False
                if not has_text:
                    return  # empty structural footer placeholder — renders nothing
            else:
                return  # date / slide-number / title / body — never branding
        if _intersect_area(r, zone) > 0:
            found.append((source, shape.name, r))

    for sh in layout.shapes:
        consider(sh, "layout")
    for sh in layout.slide_master.shapes:
        consider(sh, "master")
    return zone, found


def check(path, band=0.14, coverage=0.5):
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    flags = []
    for i, slide in enumerate(prs.slides):
        zone, branding = footer_branding(slide.slide_layout, W, H, band)
        if not branding:
            continue
        zone_area = zone[2] * zone[3]
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            r = _rect(sh)
            if not r:
                continue
            cover = _intersect_area(r, zone) / zone_area
            if cover >= coverage:
                flags.append(
                    {
                        "slide": i,
                        "picture": sh.name,
                        "zone_coverage": round(cover, 2),
                        "collides_with": [
                            {"source": s, "name": n} for (s, n, _r) in branding
                        ],
                    }
                )
                break
    return flags


def main(argv=None):
    ap = argparse.ArgumentParser(description="Doubled-logo gate for built decks.")
    ap.add_argument("pptx", help="path to the built .pptx")
    ap.add_argument(
        "--band",
        type=float,
        default=0.14,
        help="footer/logo zone height as a fraction of slide height (default 0.14)",
    )
    ap.add_argument(
        "--coverage",
        type=float,
        default=0.5,
        help="min fraction of the footer zone a picture must cover to flag (default 0.5)",
    )
    ap.add_argument("--json", action="store_true", help="emit findings as JSON")
    args = ap.parse_args(argv)

    try:
        flags = check(args.pptx, band=args.band, coverage=args.coverage)
    except Exception as e:  # noqa: BLE001 — surface any load/parse failure to the caller
        print("check_logos: error reading %s: %s" % (args.pptx, e), file=sys.stderr)
        return 2

    if args.json:
        print(json.dumps({"flagged": flags}, indent=2))
    elif not flags:
        print("check_logos: clean — no slide covers a branded footer/logo zone.")
    else:
        print(
            "check_logos: DOUBLED-LOGO RISK on %d slide(s) — a full-bleed/footer-"
            "covering picture is stacked on branding the master/layout already draws."
            % len(flags)
        )
        for f in flags:
            who = ", ".join("%s:%s" % (c["source"], c["name"]) for c in f["collides_with"])
            print(
                "  slide %d: picture '%s' covers %.0f%% of the footer zone over [%s]"
                % (f["slide"], f["picture"], f["zone_coverage"] * 100, who)
            )
        print(
            "  FIX: rebuild the picture with the logo/footer zone empty (reserve the "
            "bottom margin), or place it on a blank layout (slide_layouts[6]) so the "
            "master supplies exactly one logo. Verify by eye against the master."
        )
    return 1 if flags else 0


if __name__ == "__main__":
    sys.exit(main())
