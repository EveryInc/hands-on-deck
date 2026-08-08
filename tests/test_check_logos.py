"""Tests for check_logos.py — the doubled-logo post-build gate.

No binary fixtures: each test builds its deck with python-pptx, then drives the
gate exactly as the build workflow would (through its CLI + JSON), so the test
proves the real command flags the defect and passes the sanctioned patterns.
"""
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

SCRIPTS = Path(__file__).resolve().parents[1] / "skills" / "hands-on-deck" / "scripts"
CHECK = SCRIPTS / "check_logos.py"


def run(*args):
    return subprocess.run(
        [sys.executable, str(CHECK), *map(str, args)], capture_output=True, text=True
    )


def _brand_footer(prs, layout_idx=6):
    """Give the layout's footer placeholder real text — a rendered footer/logo that
    sits in the bottom band, exactly what a full-bleed picture would stack on.
    (python-pptx can't add graphics to a master/layout, so a branded footer
    placeholder stands in for the baked EVERY logo/footer.)"""
    layout = prs.slide_layouts[layout_idx]
    for ph in layout.placeholders:
        if "Footer" in ph.name:
            ph.text = "EVERY  ·  every.to"
            return ph
    raise AssertionError("no footer placeholder on layout %d" % layout_idx)


def _fullbleed_picture(slide, prs, img, height=None):
    """Place a full-width picture; height defaults to full-bleed (covers footer)."""
    h = height if height is not None else prs.slide_height
    slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=h)


def _png(tmp_path):
    # 2x2 red PNG — smallest valid raster the picture part will accept.
    from PIL import Image

    p = tmp_path / "bleed.png"
    Image.new("RGB", (2, 2), (200, 40, 40)).save(p)
    return p


def test_flags_fullbleed_over_branded_master(tmp_path):
    """Full-bleed picture on a layout whose master carries a logo → flagged."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    _brand_footer(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout, branded master
    _fullbleed_picture(slide, prs, _png(tmp_path))
    deck = tmp_path / "defect.pptx"
    prs.save(deck)

    r = run(deck, "--json")
    assert r.returncode == 1, r.stdout + r.stderr
    flagged = json.loads(r.stdout)["flagged"]
    assert len(flagged) == 1
    assert flagged[0]["slide"] == 0
    assert flagged[0]["zone_coverage"] >= 0.5


def test_clean_when_footer_reserved(tmp_path):
    """Same branded master, but the picture reserves the bottom margin → clean."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    _brand_footer(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Stop the image well above the footer band (band default = bottom 14%).
    _fullbleed_picture(slide, prs, _png(tmp_path), height=Inches(6.0))
    deck = tmp_path / "reserved.pptx"
    prs.save(deck)

    r = run(deck)
    assert r.returncode == 0, r.stdout + r.stderr
    assert "clean" in r.stdout


def test_clean_when_no_master_branding(tmp_path):
    """Full-bleed picture on a plain, unbranded master → sanctioned, clean."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fullbleed_picture(slide, prs, _png(tmp_path))
    deck = tmp_path / "ok.pptx"
    prs.save(deck)

    r = run(deck)
    assert r.returncode == 0, r.stdout + r.stderr
